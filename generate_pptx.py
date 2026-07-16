import sys
import subprocess
import os

# Try importing pptx, fallback to pip install without accent printing
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("Library python-pptx not found. Installing...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Khởi tạo bản trình bày
    prs = Presentation()
    
    # Thiết lập kích thước slide rộng 16:9 tiêu chuẩn
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Layout trống để dễ dàng vẽ tự do
    blank_layout = prs.slide_layouts[6]
    
    # Bảng màu chủ đạo (Indigo & Slate)
    c_dark_slate = RGBColor(15, 23, 42)
    c_indigo = RGBColor(79, 70, 229)
    c_cyan = RGBColor(6, 182, 212)
    c_text_main = RGBColor(51, 65, 85)
    c_text_white = RGBColor(255, 255, 255)
    c_bg_light = RGBColor(248, 250, 252)
    
    # Màu sắc QC
    c_green = RGBColor(16, 185, 129)
    c_amber = RGBColor(245, 158, 11)
    c_red = RGBColor(239, 68, 68)
    
    # Đường dẫn thư mục ảnh chạy train
    run_dir = "runs/segment/AI/train/runs/ai-segmentation/segmentation-yolo11n-12-07-0h"
    
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Inter'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = c_dark_slate
        
        # Đường kẻ trang trí bên dưới tiêu đề
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.2), Inches(11.83), Inches(0.04))
        line.fill.solid()
        line.fill.fore_color.rgb = c_indigo
        line.line.fill.background()

    # -------------------------------------------------------------
    # SLIDE 1: Title Slide (Dark Theme)
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    set_slide_background(s1, c_dark_slate)
    
    # Khối trang trí bên trái
    left_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = c_indigo
    left_bar.line.fill.background()
    
    # Text box tiêu đề chính
    title_box = s1.shapes.add_textbox(Inches(1.2), Inches(1.8), Inches(11.0), Inches(2.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "HỆ THỐNG TỰ ĐỘNG KIỂM TRA CHẤT LƯỢNG\nSẢN PHẨM BẰNG THỊ GIÁC MÁY TÌNH"
    p.font.name = 'Inter'
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = c_text_white
    
    # Text box tiêu đề phụ
    sub_box = s1.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(11.0), Inches(1.0))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Ứng dụng YOLO11-seg, ResNet18 Classifier & Động cơ Hỏi đáp VQA đa ngôn ngữ"
    p_sub.font.name = 'Inter'
    p_sub.font.size = Pt(20)
    p_sub.font.italic = True
    p_sub.font.color.rgb = c_cyan
    
    info_box = s1.shapes.add_textbox(Inches(1.2), Inches(5.8), Inches(11.0), Inches(0.8))
    p_info = info_box.text_frame.paragraphs[0]
    p_info.text = "Báo cáo Đồ án tốt nghiệp / Nghiên cứu khoa học chuyên ngành Thị giác Máy tính"
    p_info.font.name = 'Inter'
    p_info.font.size = Pt(14)
    p_info.font.color.rgb = RGBColor(148, 163, 184)

    # -------------------------------------------------------------
    # SLIDE 2: Đặt vấn đề & Mục tiêu
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    set_slide_background(s2, c_bg_light)
    add_slide_header(s2, "1. Đặt Vấn Đề & Mục Tiêu Dự Án")
    
    # Cột bên trái: Thực trạng
    col1 = s2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "THỰC TRẠNG & THÁCH THỨC"
    p1.font.name = 'Inter'
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = c_indigo
    
    bullets1 = [
        "Kiểm tra chất lượng (QC) thủ công bằng mắt dễ gây sai sót do sự mệt mỏi của con người.",
        "Thiếu số liệu định lượng về kích thước lỗi, diện tích phủ của lỗi và vị trí lỗi cụ thể trên sản phẩm.",
        "Tốc độ kiểm định chậm, gây tắc nghẽn dây chuyền sản xuất hàng loạt.",
        "Khó lưu trữ lịch sử để cải tiến quy trình QC."
    ]
    for b in bullets1:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Inter'
        p.font.size = Pt(15)
        p.font.color.rgb = c_text_main
        p.space_before = Pt(14)

    # Cột bên phải: Mục tiêu
    col2 = s2.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.8))
    tf2 = col2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = "MỤC TIÊU CỦA GIẢI PHÁP AI"
    p2.font.name = 'Inter'
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = c_cyan
    
    bullets2 = [
        "Tự động hóa hoàn toàn quy trình phát hiện và khoanh vùng lỗi bề mặt thời gian thực (Real-time).",
        "Trích xuất số liệu hình học chính xác diện tích, phân vùng lưới 9 khu vực và tính điểm nghiêm trọng.",
        "Xây dựng Dashboard trực quan hóa và bộ giả lập Conveyor Live cho sản xuất thực tế.",
        "Tích hợp công nghệ Chatbot VQA hỗ trợ cả tiếng Anh lẫn tiếng Việt giúp tương tác dễ dàng với hệ thống."
    ]
    for b in bullets2:
        p = tf2.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Inter'
        p.font.size = Pt(15)
        p.font.color.rgb = c_text_main
        p.space_before = Pt(14)

    # -------------------------------------------------------------
    # SLIDE 3: Kiến trúc hệ thống (Sơ đồ khối vẽ bằng code)
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    set_slide_background(s3, c_bg_light)
    add_slide_header(s3, "2. Kiến Trúc Luồng Xử Lý Hệ Thống")
    
    # Vẽ các Block node hệ thống
    x_positions = [Inches(0.75), Inches(3.25), Inches(5.75), Inches(8.25), Inches(10.75)]
    titles = ["💻 Client\n(HTML/JS)", "⚡ FastAPI\nBackend", "🔮 Model 1\nYOLO11-seg", "🔍 Model 2\nResNet18", "📊 Logic\nQC & VQA"]
    colors = [c_dark_slate, c_dark_slate, c_indigo, c_cyan, c_green]
    
    # Vẽ các khối nút
    for i in range(5):
        node = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_positions[i], Inches(2.2), Inches(1.8), Inches(1.4))
        node.fill.solid()
        node.fill.fore_color.rgb = colors[i]
        node.line.color.rgb = colors[i]
        
        tf = node.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = titles[i]
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = c_text_white
        p.alignment = PP_ALIGN.CENTER
        
        # Thêm mũi tên kết nối
        if i < 4:
            arrow = s3.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x_positions[i] + Inches(1.85), Inches(2.75), Inches(0.6), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = c_indigo
            arrow.line.fill.background()

    # Thêm text mô tả chi tiết luồng bên dưới sơ đồ
    desc_box = s3.shapes.add_textbox(Inches(0.75), Inches(4.2), Inches(11.83), Inches(2.5))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p_desc = tf_desc.paragraphs[0]
    p_desc.text = "CHI TIẾT LUỒNG TRUYỀN DỮ LIỆU:"
    p_desc.font.name = 'Inter'
    p_desc.font.size = Pt(16)
    p_desc.font.bold = True
    p_desc.font.color.rgb = c_dark_slate
    
    steps = [
        "Bước 1: Client tải lên ảnh hoặc video từ giao diện Dashboard thông minh.",
        "Bước 2: FastAPI tiếp nhận yêu cầu, định tuyến file và chuyển tiếp sang module AI xử lý.",
        "Bước 3: Model YOLO11-seg xác định vùng lỗi và xuất ra mặt nạ đa giác (Segmentation Mask).",
        "Bước 4: Cắt vùng ảnh lỗi (ROI), đưa sang ResNet18 Classifier để nhận diện nhãn lỗi chi tiết (Dent, Scratch,...).",
        "Bước 5: Trích xuất diện tích toán học, phân lưới vị trí, tính điểm nghiêm trọng và tạo chatbot VQA phản hồi."
    ]
    for step in steps:
        p = tf_desc.add_paragraph()
        p.text = step
        p.font.name = 'Inter'
        p.font.size = Pt(13)
        p.font.color.rgb = c_text_main
        p.space_before = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 4: Bộ dữ liệu MVTec-AD (Có ảnh minh họa thực tế)
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    set_slide_background(s4, c_bg_light)
    add_slide_header(s4, "3. Tập Dữ Liệu Huấn Luyện MVTec-AD")
    
    # Cột chữ bên trái
    col1 = s4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "TẬP DỮ LIỆU & TIỀN XỬ LÝ"
    p1.font.name = 'Inter'
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = c_indigo
    
    bullets = [
        "MVTec Anomaly Detection (MVTec-AD) là bộ dữ liệu chuẩn mực để đánh giá các thuật toán kiểm tra lỗi bề mặt công nghiệp.",
        "Bao gồm các nhóm sản phẩm cố định, chia làm 2 tập: Normal (Huấn luyện) và Anomaly (Xác thực/Kiểm thử).",
        "Kích thước ảnh được resize chuẩn hóa về 640x640 pixel.",
        "Mặt nạ nhãn lỗi được convert đồng bộ sang định dạng YOLO Segment phục vụ quá trình huấn luyện."
    ]
    for b in bullets:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.color.rgb = c_text_main
        p.space_before = Pt(12)
        
    # Cột ảnh bên phải: Load ảnh labels.jpg từ thư mục train
    img_labels = os.path.join(run_dir, "labels.jpg")
    if os.path.exists(img_labels):
        s4.shapes.add_picture(img_labels, Inches(6.6), Inches(1.8), width=Inches(6.0), height=Inches(4.5))
    else:
        # Hộp thoại dự phòng nếu thiếu ảnh
        placeholder = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.8), Inches(6.0), Inches(4.5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(226, 232, 240)
        placeholder.line.color.rgb = c_indigo
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.text = "[Hình ảnh nhãn huấn luyện thực tế]\n\n(Tệp labels.jpg hiển thị các nhãn lớp và đa giác khoanh vùng lỗi MVTec-AD)"
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.color.rgb = c_indigo
        p.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 5: Tăng cường dữ liệu (Augmentation)
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    set_slide_background(s5, c_bg_light)
    add_slide_header(s5, "4. Kỹ Thuật Tăng Cường Dữ Liệu (Augmentation)")
    
    # Thiết kế 4 khối biểu diễn 4 phương pháp Augment
    x_positions = [Inches(0.75), Inches(3.75), Inches(6.75), Inches(9.75)]
    titles = ["🔄 XOAY & LẬT", "☀️ ĐỘ SÁNG", "🔊 NHIỄU GAUSS", "🔎 THU PHÓNG"]
    descs = [
        "Lập ảnh ngẫu nhiên theo chiều dọc/ngang (Horizontal & Vertical Flip) và xoay tự do.\n\nGiúp mô hình nhận diện khuyết tật không phụ thuộc vào góc quay của sản phẩm trên băng tải.",
        "Thay đổi ngẫu nhiên độ sáng, độ tương phản (Brightness/Contrast).\n\nGiúp mô hình ổn định ngay cả khi ánh sáng nhà xưởng thay đổi đột ngột giữa ca ngày và đêm.",
        "Thêm nhiễu Gaussian ngẫu nhiên ở các cấp độ hạt mịn khác nhau.\n\nMô phỏng nhiễu phần cứng từ camera công nghiệp trong môi trường rung động của nhà máy.",
        "Cắt ngẫu nhiên và thu phóng (Random Crop & Scale).\n\nTập trung huấn luyện mô hình nhận diện lỗi ở nhiều tỷ lệ khoảng cách camera xa gần khác nhau."
    ]
    
    for i in range(4):
        card = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_positions[i], Inches(2.0), Inches(2.83), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = c_dark_slate
        card.line.color.rgb = c_cyan
        
        tf = card.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = titles[i]
        p.font.name = 'Inter'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = c_cyan
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(14)
        
        p2 = tf.add_paragraph()
        p2.text = descs[i]
        p2.font.name = 'Inter'
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(226, 232, 240)
        p2.alignment = PP_ALIGN.LEFT

    # -------------------------------------------------------------
    # SLIDE 6: Phát hiện & Phân vùng lỗi (Model 1: YOLO11-seg)
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_layout)
    set_slide_background(s6, c_bg_light)
    add_slide_header(s6, "5. Phân Vùng Lỗi Bề Mặt Bằng YOLO11-seg")
    
    # Cột chữ bên trái
    col1 = s6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "KẾT QUẢ ĐÀO TẠO MÔ HÌNH"
    p1.font.name = 'Inter'
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = c_indigo
    
    bullets = [
        "YOLO11-seg làm mô hình phát hiện và tạo mặt nạ phân vùng lỗi (mask) chính xác tới cấp độ pixel.",
        "Quá trình huấn luyện tối ưu hóa hai chỉ số chính: mAP (Mean Average Precision) và IoU (Intersection over Union).",
        "Khống chế tỷ lệ bỏ sót khuyết tật nghiêm trọng (False Negative) dưới mức an toàn 1-2%.",
        "Trọng số huấn luyện tốt nhất được xuất thành best.pt để làm lõi suy luận cho API."
    ]
    for b in bullets:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.color.rgb = c_text_main
        p.space_before = Pt(12)
        
    # Cột ảnh bên phải: Kết quả chạy train (results.png)
    img_results = os.path.join(run_dir, "results.png")
    if os.path.exists(img_results):
        s6.shapes.add_picture(img_results, Inches(6.6), Inches(1.8), width=Inches(6.0), height=Inches(4.5))
    else:
        # Hộp thoại dự phòng nếu thiếu ảnh
        placeholder = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.8), Inches(6.0), Inches(4.5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(226, 232, 240)
        placeholder.line.color.rgb = c_indigo
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.text = "[Biểu đồ kết quả huấn luyện]\n\n(Tệp results.png hiển thị các đường cong Loss và mAP trong quá trình Train/Val)"
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.color.rgb = c_indigo
        p.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 7: Phân loại loại lỗi chi tiết (Model 2: ResNet18)
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_layout)
    set_slide_background(s7, c_bg_light)
    add_slide_header(s7, "6. Phân Loại Loại Lỗi Bằng ResNet18 Classifier")
    
    # Cột chữ bên trái
    col1 = s7.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "NHẬN DIỆN CHI TIẾT TỪNG NHÓM LỖI"
    p1.font.name = 'Inter'
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = c_indigo
    
    bullets = [
        "Mô hình ResNet18 Classifier chuyên biệt chạy nối tiếp sau khi YOLO11-seg phát hiện vùng khuyết tật.",
        "Cắt vùng ảnh chứa lỗi (ROI), mở rộng 15% biên để tránh mất đặc trưng góc cạnh, đưa vào phân loại chi tiết.",
        "Ứng dụng file cấu hình JSON ràng buộc logic: chỉ cho phép phân loại lỗi thuộc danh mục vật lý của nhóm sản phẩm cha.",
        "Ví dụ minh họa bên phải hiển thị ảnh dự đoán kiểm định (val_batch0_pred) thực tế của mô hình."
    ]
    for b in bullets:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.color.rgb = c_text_main
        p.space_before = Pt(12)
        
    # Cột ảnh bên phải: Ảnh dự đoán thực tế (val_batch0_pred.jpg)
    img_preds = os.path.join(run_dir, "val_batch0_pred.jpg")
    if os.path.exists(img_preds):
        s7.shapes.add_picture(img_preds, Inches(6.6), Inches(1.8), width=Inches(6.0), height=Inches(4.5))
    else:
        # Hộp thoại dự phòng nếu thiếu ảnh
        placeholder = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.8), Inches(6.0), Inches(4.5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(226, 232, 240)
        placeholder.line.color.rgb = c_indigo
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.text = "[Ảnh dự đoán khuyết tật thực tế]\n\n(Tệp val_batch0_pred.jpg hiển thị các nhãn lỗi đè lên vùng ảnh thực tế)"
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.color.rgb = c_indigo
        p.alignment = PP_ALIGN.CENTER

    # -------------------------------------------------------------
    # SLIDE 8: Trích xuất đặc trưng & Luật QC (Có biểu đồ vẽ bằng shape)
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_layout)
    set_slide_background(s8, c_bg_light)
    add_slide_header(s8, "7. Thuật Toán Trích Xuất & Quyết Định QC")
    
    # Cột bên trái: Trực quan hóa trọng số tính điểm bằng shape
    left_title = s8.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(0.5))
    left_title.text_frame.paragraphs[0].text = "PHÂN BỔ TRỌNG SỐ SEVERITY (0-100)"
    left_title.text_frame.paragraphs[0].font.name = 'Inter'
    left_title.text_frame.paragraphs[0].font.size = Pt(18)
    left_title.text_frame.paragraphs[0].font.bold = True
    left_title.text_frame.paragraphs[0].font.color.rgb = c_indigo
    
    labels = ["Diện tích lỗi (Shoelace)", "Vị trí lỗi (Lưới 9 vùng)", "Độ tin cậy mô hình AI"]
    weights = [0.4, 0.3, 0.3]
    fills = [c_indigo, c_cyan, c_green]
    
    for i in range(3):
        # Nhãn
        lbl_box = s8.shapes.add_textbox(Inches(0.75), Inches(2.4 + i*1.0), Inches(5.5), Inches(0.4))
        p = lbl_box.text_frame.paragraphs[0]
        p.text = f"{labels[i]} ({int(weights[i]*100)}%)"
        p.font.name = 'Inter'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = c_text_main
        
        # Thanh nền
        bg_bar = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.8 + i*1.0), Inches(5.0), Inches(0.2))
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = RGBColor(226, 232, 240)
        bg_bar.line.fill.background()
        
        # Thanh đầy
        fill_bar = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.8 + i*1.0), Inches(5.0 * weights[i]), Inches(0.2))
        fill_bar.fill.solid()
        fill_bar.fill.fore_color.rgb = fills[i]
        fill_bar.line.fill.background()
        
    # Công thức toán học
    f_box = s8.shapes.add_textbox(Inches(0.75), Inches(5.5), Inches(5.0), Inches(1.0))
    tf_f = f_box.text_frame
    tf_f.word_wrap = True
    p_f = tf_f.paragraphs[0]
    p_f.text = "Công thức Severity Score:\nSeverity = (Area * 0.4) + (Pos * 0.3) + (Conf * 0.3)"
    p_f.font.name = 'Courier New'
    p_f.font.size = Pt(13)
    p_f.font.bold = True
    p_f.font.color.rgb = c_indigo
    
    # Cột bên phải: Quyết định QC (Có vẽ thẻ màu bằng Shape)
    col2_title = s8.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(0.5))
    col2_title.text_frame.paragraphs[0].text = "HỆ THỐNG QUYẾT ĐỊNH QC CHUẨN"
    col2_title.text_frame.paragraphs[0].font.name = 'Inter'
    col2_title.text_frame.paragraphs[0].font.size = Pt(18)
    col2_title.text_frame.paragraphs[0].font.bold = True
    col2_title.text_frame.paragraphs[0].font.color.rgb = c_cyan
    
    decisions = [
        ("PASS (ĐẠT)", "Sản phẩm không có lỗi hoặc toàn bộ lỗi được phát hiện chỉ có mức độ nghiêm trọng thấp (Severity < 25).", c_green),
        ("FLAG (CẢNH BÁO)", "Phát hiện có từ 1 đến 2 lỗi ở mức độ trung bình (Severity: 25 - 50) ➔ Chuyển qua luồng kiểm định lại thủ công.", c_amber),
        ("REJECT (LỖI)", "Có từ 3 lỗi Medium trở lên, hoặc xuất hiện bất kỳ lỗi nào ở mức độ cao (Severity > 50) ➔ Huỷ sản phẩm.", c_red)
    ]
    for i, (title, desc, color) in enumerate(decisions):
        card = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.4 + i*1.5), Inches(5.8), Inches(1.2))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(241, 245, 249)
        card.line.color.rgb = color
        card.line.width = Pt(2)
        
        tf = card.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Inter'
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = c_text_main
        p2.space_before = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 9: Động cơ Hỏi Đáp VQA (Có sơ đồ luồng shape)
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_layout)
    set_slide_background(s9, c_bg_light)
    add_slide_header(s9, "8. Động Cơ Hỏi Đáp Ngữ Cảnh Đa Ngôn Ngữ (VQA)")
    
    # Cột chữ bên trái
    col1 = s9.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "VQA CHATBOT ĐA NGÔN NGỮ"
    p1.font.name = 'Inter'
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = c_indigo
    
    bullets = [
        "Người vận hành tương tác trực tiếp bằng câu hỏi ngôn ngữ tự nhiên để tra cứu thông tin lỗi nhanh chóng.",
        "Bộ quét ngôn ngữ tự động phát hiện câu hỏi tiếng Việt dựa trên sự xuất hiện của các từ khóa đặc trưng tiếng Việt.",
        "Phản hồi tiếng Việt động trực tiếp từ kết quả phân tích chất lượng của sản phẩm (loại lỗi, vị trí, kết luận, khuyến cáo).",
        "Có lưu trữ context cục bộ tại frontend để tự động chuyển tiếp dữ liệu phân tích sang VQA chat bubble."
    ]
    for b in bullets:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.color.rgb = c_text_main
        p.space_before = Pt(12)
        
    # Cột vẽ sơ đồ phân cấp VQA bên phải
    col2_title = s9.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(0.5))
    col2_title.text_frame.paragraphs[0].text = "QUY TRÌNH PHÂN CẤP ƯU TIÊN 3 TẦNG"
    col2_title.text_frame.paragraphs[0].font.name = 'Inter'
    col2_title.text_frame.paragraphs[0].font.size = Pt(18)
    col2_title.text_frame.paragraphs[0].font.bold = True
    col2_title.text_frame.paragraphs[0].font.color.rgb = c_cyan
    
    vqa_levels = [
        ("TẦNG 1: CONTEXT-AWARE", "Trích xuất nhanh số liệu từ báo cáo QC. Phản hồi tiếng Việt/Anh tức thì.", c_cyan),
        ("TẦNG 2: DEEP LEARNING (ViLT)", "Gọi model dandelin/vilt-b32 để phân tích trực quan tổng thể.", c_indigo),
        ("TẦNG 3: KEYWORD FALLBACK", "Đối khớp từ khóa thô khi mất kết nối mạng hoặc không có ảnh.", c_dark_slate)
    ]
    for i, (title, desc, color) in enumerate(vqa_levels):
        box = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.4 + i*1.5), Inches(5.8), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Inter'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = c_text_white
        
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Inter'
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(226, 232, 240)
        p2.space_before = Pt(4)
        
        # Mũi tên đi xuống giữa các tầng
        if i < 2:
            arrow = s9.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9.5), Inches(3.55 + i*1.5), Inches(0.4), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = c_indigo
            arrow.line.fill.background()

    # -------------------------------------------------------------
    # SLIDE 10: Kết luận & Hướng phát triển (Có ảnh Matrix thực tế)
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_layout)
    set_slide_background(s10, c_dark_slate)
    
    # Khối bên trái trang trí
    left_bar = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = c_cyan
    left_bar.line.fill.background()
    
    # Cột chữ bên trái
    col1 = s10.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(5.2), Inches(5.0))
    tf1 = col1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = "KẾT QUẢ ĐẠT ĐƯỢC"
    p1.font.name = 'Inter'
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = c_cyan
    
    bullets = [
        "Xây dựng thành công hệ thống tích hợp mượt mà giữa AI Vision và Dashboard quản lý.",
        "Tỉ lệ bỏ sót lỗi nghiêm trọng rất thấp (< 1-2%), đáp ứng yêu cầu chất lượng công nghiệp.",
        "Dashboard trực quan hóa cao cấp, thao tác đơn giản và hỗ trợ VQA phản hồi tiếng Việt cực nhạy."
    ]
    for b in bullets:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_before = Pt(16)
        
    p_future = tf1.add_paragraph()
    p_future.text = "HƯỚNG PHÁT TRIỂN TIẾP THEO"
    p_future.font.name = 'Inter'
    p_future.font.size = Pt(24)
    p_future.font.bold = True
    p_future.font.color.rgb = c_text_white
    p_future.space_before = Pt(30)
    
    bullets2 = [
        "Tối ưu hóa tốc độ suy luận bằng TensorRT để chạy ổn định trên thiết bị nhúng Nvidia Jetson cạnh băng chuyền.",
        "Tích hợp kết nối trực tiếp với các thiết bị vật lý công nghiệp PLC, cánh tay robot để tự động gạt sản phẩm lỗi."
    ]
    for b in bullets2:
        p = tf1.add_paragraph()
        p.text = "• " + b
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_before = Pt(12)

    # Cột ảnh bên phải: Confusion Matrix (confusion_matrix.png)
    img_matrix = os.path.join(run_dir, "confusion_matrix.png")
    if os.path.exists(img_matrix):
        s10.shapes.add_picture(img_matrix, Inches(6.8), Inches(1.5), width=Inches(5.5), height=Inches(4.5))
    else:
        # Hộp thoại dự phòng nếu thiếu ảnh
        placeholder = s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.5), Inches(5.5), Inches(4.5))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = RGBColor(30, 41, 59)
        placeholder.line.color.rgb = c_cyan
        tf = placeholder.text_frame
        p = tf.paragraphs[0]
        p.text = "[Ma trận nhầm lẫn - Confusion Matrix]\n\n(Tệp confusion_matrix.png hiển thị hiệu suất phân loại thực tế của mô hình)"
        p.font.name = 'Inter'
        p.font.size = Pt(14)
        p.font.color.rgb = c_cyan
        p.alignment = PP_ALIGN.CENTER

    # Lưu tệp PowerPoint
    output_filename = "presentation_ai_qc_with_diagrams.pptx"
    prs.save(output_filename)
    print("Presentation saved successfully: " + os.path.abspath(output_filename))

if __name__ == "__main__":
    create_presentation()
